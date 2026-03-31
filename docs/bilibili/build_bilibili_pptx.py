from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


PAGE_WIDTH = Inches(13.333)
PAGE_HEIGHT = Inches(7.5)

COLOR_NAVY = RGBColor(18, 45, 90)
COLOR_BLUE = RGBColor(43, 108, 176)
COLOR_TEAL = RGBColor(45, 146, 166)
COLOR_GOLD = RGBColor(214, 158, 46)
COLOR_LIGHT = RGBColor(245, 248, 252)
COLOR_TEXT = RGBColor(35, 41, 52)
COLOR_MUTED = RGBColor(94, 104, 121)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_LINE = RGBColor(220, 228, 239)


@dataclass(frozen=True)
class SlideSpec:
    layout: str
    title: str
    bullets: list[str] | None = None
    subtitle: str | None = None
    code: str | None = None
    image: str | None = None
    image_caption: str | None = None
    left_title: str | None = None
    left_bullets: list[str] | None = None
    right_title: str | None = None
    right_bullets: list[str] | None = None


@dataclass(frozen=True)
class LectureDeck:
    slug: str
    title: str
    subtitle: str
    accent: RGBColor
    slides: list[SlideSpec]


REPO_ROOT = Path(__file__).resolve().parents[2]
BILIBILI_DIR = Path(__file__).resolve().parent


def apply_background(slide, accent: RGBColor) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = COLOR_WHITE

    top_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PAGE_WIDTH, Inches(0.42))
    top_band.fill.solid()
    top_band.fill.fore_color.rgb = accent
    top_band.line.fill.background()

    bottom_band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, PAGE_HEIGHT - Inches(0.18), PAGE_WIDTH, Inches(0.18)
    )
    bottom_band.fill.solid()
    bottom_band.fill.fore_color.rgb = COLOR_NAVY
    bottom_band.line.fill.background()


def add_footer(slide, text: str) -> None:
    textbox = slide.shapes.add_textbox(Inches(0.45), PAGE_HEIGHT - Inches(0.42), Inches(12.2), Inches(0.18))
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(9)
    paragraph.font.color.rgb = COLOR_MUTED
    paragraph.alignment = PP_ALIGN.RIGHT


def set_text_style(paragraph, size: int, color: RGBColor = COLOR_TEXT, bold: bool = False) -> None:
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.font.name = "Consolas"


def add_title(slide, title: str, subtitle: str | None = None, accent: RGBColor = COLOR_BLUE) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.72), Inches(8.9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    paragraph = title_frame.paragraphs[0]
    paragraph.text = title
    set_text_style(paragraph, 24, COLOR_NAVY, True)

    title_rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.68), Inches(1.55), Inches(1.4), Inches(0.06))
    title_rule.fill.solid()
    title_rule.fill.fore_color.rgb = accent
    title_rule.line.fill.background()

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.68), Inches(1.75), Inches(11.6), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        paragraph = subtitle_frame.paragraphs[0]
        paragraph.text = subtitle
        set_text_style(paragraph, 12, COLOR_MUTED, False)


def add_bullet_list(
    slide,
    bullets: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 20,
) -> None:
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.space_after = Pt(10)
        paragraph.bullet = True
        set_text_style(paragraph, font_size)


def add_code_block(slide, code: str, left: float, top: float, width: float, height: float) -> None:
    code_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    code_panel.fill.solid()
    code_panel.fill.fore_color.rgb = COLOR_LIGHT
    code_panel.line.color.rgb = COLOR_LINE

    textbox = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), height - Inches(0.22))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = text_frame.paragraphs[0]
    paragraph.text = code
    paragraph.font.name = "Consolas"
    paragraph.font.size = Pt(15)
    paragraph.font.color.rgb = COLOR_TEXT


def add_image_card(slide, image_path: Path, caption: str | None, left: float, top: float, width: float, height: float) -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLOR_WHITE
    panel.line.color.rgb = COLOR_LINE

    available_left = left + Inches(0.1)
    available_top = top + Inches(0.1)
    available_width = width - Inches(0.2)
    available_height = height - Inches(0.36 if caption else 0.2)
    slide.shapes.add_picture(str(image_path), available_left, available_top, width=available_width, height=available_height)

    if caption:
        caption_box = slide.shapes.add_textbox(left + Inches(0.14), top + height - Inches(0.28), width - Inches(0.28), Inches(0.18))
        paragraph = caption_box.text_frame.paragraphs[0]
        paragraph.text = caption
        set_text_style(paragraph, 10, COLOR_MUTED, False)
        paragraph.alignment = PP_ALIGN.CENTER


def add_section_tag(slide, text: str, accent: RGBColor) -> None:
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.8), Inches(0.63), Inches(1.8), Inches(0.35))
    pill.fill.solid()
    pill.fill.fore_color.rgb = accent
    pill.line.fill.background()
    paragraph = slide.shapes.add_textbox(Inches(10.8), Inches(0.67), Inches(1.8), Inches(0.25)).text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    set_text_style(paragraph, 10, COLOR_WHITE, True)


def build_title_slide(prs: Presentation, deck: LectureDeck) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, deck.accent)

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(8.4), Inches(1.0))
    title_paragraph = title_box.text_frame.paragraphs[0]
    title_paragraph.text = deck.title
    set_text_style(title_paragraph, 28, COLOR_NAVY, True)

    subtitle_box = slide.shapes.add_textbox(Inches(0.78), Inches(2.1), Inches(9.0), Inches(0.9))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_paragraph = subtitle_frame.paragraphs[0]
    subtitle_paragraph.text = deck.subtitle
    set_text_style(subtitle_paragraph, 16, COLOR_MUTED, False)

    highlight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.25), Inches(3.4), Inches(0.65))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = deck.accent
    highlight.line.fill.background()
    paragraph = slide.shapes.add_textbox(Inches(0.98), Inches(3.42), Inches(3.0), Inches(0.25)).text_frame.paragraphs[0]
    paragraph.text = "项目导向 · 代码导向 · 可录制初稿"
    set_text_style(paragraph, 13, COLOR_WHITE, True)

    info_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.95), Inches(1.2), Inches(3.45), Inches(4.2))
    info_box.fill.solid()
    info_box.fill.fore_color.rgb = COLOR_LIGHT
    info_box.line.color.rgb = COLOR_LINE

    info_text = slide.shapes.add_textbox(Inches(9.18), Inches(1.52), Inches(3.0), Inches(3.6))
    info_frame = info_text.text_frame
    info_frame.word_wrap = True

    items = [
        "适合：源码讲解型 B 站课程",
        "素材来源：当前仓库 `src/`、`pupper/`、`sim/`",
        "讲法：架构图 -> 源码 -> 演示 -> 作业",
        "用途：正式课件结构稿，可继续美化",
    ]
    for index, item in enumerate(items):
        paragraph = info_frame.paragraphs[0] if index == 0 else info_frame.add_paragraph()
        paragraph.text = item
        paragraph.bullet = True
        paragraph.space_after = Pt(10)
        set_text_style(paragraph, 14)

    add_footer(slide, f"{deck.slug} · StanfordQuadruped B 站课程")


def build_bullets_slide(prs: Presentation, deck: LectureDeck, spec: SlideSpec, slide_index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, deck.accent)
    add_title(slide, spec.title, spec.subtitle, deck.accent)
    add_section_tag(slide, f"{slide_index:02d}", deck.accent)
    add_bullet_list(slide, spec.bullets or [], Inches(0.88), Inches(2.15), Inches(11.2), Inches(4.5))
    add_footer(slide, f"{deck.slug} · 第 {slide_index} 页")


def build_image_slide(prs: Presentation, deck: LectureDeck, spec: SlideSpec, slide_index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, deck.accent)
    add_title(slide, spec.title, spec.subtitle, deck.accent)
    add_section_tag(slide, f"{slide_index:02d}", deck.accent)
    image_path = (REPO_ROOT / spec.image).resolve()
    add_bullet_list(slide, spec.bullets or [], Inches(0.85), Inches(2.1), Inches(4.4), Inches(4.5), font_size=18)
    add_image_card(slide, image_path, spec.image_caption, Inches(5.45), Inches(1.95), Inches(7.1), Inches(4.9))
    add_footer(slide, f"{deck.slug} · 第 {slide_index} 页")


def build_code_slide(prs: Presentation, deck: LectureDeck, spec: SlideSpec, slide_index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, deck.accent)
    add_title(slide, spec.title, spec.subtitle, deck.accent)
    add_section_tag(slide, f"{slide_index:02d}", deck.accent)
    if spec.bullets:
        add_bullet_list(slide, spec.bullets, Inches(0.82), Inches(2.0), Inches(4.3), Inches(4.8), font_size=17)
        add_code_block(slide, spec.code or "", Inches(5.25), Inches(1.95), Inches(7.1), Inches(4.95))
    else:
        add_code_block(slide, spec.code or "", Inches(0.85), Inches(2.0), Inches(11.6), Inches(4.95))
    add_footer(slide, f"{deck.slug} · 第 {slide_index} 页")


def build_two_col_slide(prs: Presentation, deck: LectureDeck, spec: SlideSpec, slide_index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, deck.accent)
    add_title(slide, spec.title, spec.subtitle, deck.accent)
    add_section_tag(slide, f"{slide_index:02d}", deck.accent)

    left_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(2.0), Inches(5.55), Inches(4.95))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = COLOR_LIGHT
    left_panel.line.color.rgb = COLOR_LINE
    right_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.55), Inches(2.0), Inches(5.95), Inches(4.95))
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = COLOR_LIGHT
    right_panel.line.color.rgb = COLOR_LINE

    left_title = slide.shapes.add_textbox(Inches(1.05), Inches(2.25), Inches(4.9), Inches(0.3))
    paragraph = left_title.text_frame.paragraphs[0]
    paragraph.text = spec.left_title or ""
    set_text_style(paragraph, 16, deck.accent, True)
    add_bullet_list(slide, spec.left_bullets or [], Inches(1.0), Inches(2.62), Inches(4.9), Inches(3.9), font_size=15)

    right_title = slide.shapes.add_textbox(Inches(6.78), Inches(2.25), Inches(5.2), Inches(0.3))
    paragraph = right_title.text_frame.paragraphs[0]
    paragraph.text = spec.right_title or ""
    set_text_style(paragraph, 16, deck.accent, True)
    add_bullet_list(slide, spec.right_bullets or [], Inches(6.74), Inches(2.62), Inches(5.1), Inches(3.9), font_size=15)
    add_footer(slide, f"{deck.slug} · 第 {slide_index} 页")


def create_deck(deck: LectureDeck) -> Path:
    presentation = Presentation()
    presentation.slide_width = PAGE_WIDTH
    presentation.slide_height = PAGE_HEIGHT

    build_title_slide(presentation, deck)
    for slide_index, spec in enumerate(deck.slides, start=1):
        if spec.layout == "bullets":
            build_bullets_slide(presentation, deck, spec, slide_index)
        elif spec.layout == "image":
            build_image_slide(presentation, deck, spec, slide_index)
        elif spec.layout == "code":
            build_code_slide(presentation, deck, spec, slide_index)
        elif spec.layout == "two_col":
            build_two_col_slide(presentation, deck, spec, slide_index)
        else:
            raise ValueError(f"Unsupported layout `{spec.layout}`.")

    output_path = BILIBILI_DIR / f"{deck.slug}_courseware.pptx"
    presentation.save(output_path)
    return output_path


def build_lectures() -> list[LectureDeck]:
    return [
        LectureDeck(
            slug="lecture01_project_overview",
            title="第 1 讲：项目全景、入口与控制主循环",
            subtitle="先把仓库地图、控制主链和三类核心数据对象讲清楚。",
            accent=COLOR_BLUE,
            slides=[
                SlideSpec(
                    layout="bullets",
                    title="本讲目标",
                    bullets=[
                        "判断项目控制范式：gait + 足端规划 + IK，而不是 MPC / WBC。",
                        "认清仓库四层结构：入口层、控制层、模型硬件层、仿真桥接层。",
                        "建立 `Command / State / Configuration` 三对象心智模型。",
                        "把 `run_robot.py` 主循环讲成一条可复述的运行链。",
                    ],
                ),
                SlideSpec(
                    layout="image",
                    title="仓库四层结构",
                    bullets=[
                        "入口层：`run_robot.py`、`sim/run_*`。",
                        "控制层：`src/Controller.py`、`Gaits.py`、`StanceController.py`、`SwingLegController.py`。",
                        "模型与硬件层：`pupper/Config.py`、`Kinematics.py`、`HardwareInterface.py`。",
                        "仿真桥接层：MuJoCo 构建、观测、执行器与任务调度。",
                    ],
                    image="docs/mujoco_quadruped_mastery_tutorial/assets/project_structure_map.png",
                    image_caption="项目结构图：最适合第 1 讲用来建立整体地图。",
                ),
                SlideSpec(
                    layout="two_col",
                    title="两个入口分别承担什么角色",
                    left_title="`run_robot.py`",
                    left_bullets=[
                        "原始实机入口。",
                        "串联手柄、IMU、控制器和舵机后端。",
                        "强调真实运行链和部署依赖。",
                    ],
                    right_title="`sim/run_floating_base.py`",
                    right_bullets=[
                        "当前最佳教学入口。",
                        "共享同一控制器核心，但更容易观察状态。",
                        "适合讲闭环、调参与回归实验。",
                    ],
                ),
                SlideSpec(
                    layout="bullets",
                    title="三类核心数据对象",
                    bullets=[
                        "`Command`：这一拍想让机器人做什么。",
                        "`State`：控制器当前维护的足端、关节、姿态和观测状态。",
                        "`Configuration`：步态、几何、输入限制和仿真参数集合。",
                        "建议观众牢记：输入、状态、参数三者分离是工程可维护性的基础。",
                    ],
                ),
                SlideSpec(
                    layout="code",
                    title="一拍控制循环伪代码",
                    bullets=[
                        "主入口负责 orchestration，不负责复杂控制算法。",
                        "真正的运动逻辑集中在 `Controller.run()`。",
                        "输出统一收敛到 `state.joint_angles`。",
                    ],
                    code=(
                        "config = Configuration()\n"
                        "hardware = HardwareInterface()\n"
                        "controller = Controller(config, four_legs_inverse_kinematics)\n"
                        "state = State()\n"
                        "joystick = JoystickInterface(config)\n\n"
                        "while active:\n"
                        "    command = joystick.get_command(state)\n"
                        "    state.quat_orientation = imu.read_orientation() or [1, 0, 0, 0]\n"
                        "    controller.run(state, command)\n"
                        "    hardware.set_actuator_postions(state.joint_angles)\n"
                    ),
                ),
                SlideSpec(
                    layout="image",
                    title="为什么先从 MuJoCo 入口学",
                    bullets=[
                        "不依赖真实硬件，课程录制更稳定。",
                        "能看到机身与地面的真实互动。",
                        "能直接读取姿态、速度、触地和关节状态。",
                        "更适合讲“共享控制核心 + 可替换后端”。",
                    ],
                    image="docs/mujoco_quadruped_mastery_tutorial/assets/entrypoints_learning_path.png",
                    image_caption="推荐学习路线：先仿真，后实机。",
                ),
                SlideSpec(
                    layout="code",
                    title="建议现场演示命令",
                    code=(
                        "python sim/run_fixed_base.py --mode rest --duration 10\n"
                        "python sim/run_fixed_base.py --mode trot --duration 10\n"
                        "python sim/run_floating_base.py --mode rest --duration 10\n\n"
                        "# 环境缺依赖时可先安装\n"
                        "pip install numpy transforms3d mujoco\n"
                    ),
                    bullets=[
                        "第 1 讲演示目标不是调参，而是让观众看到入口与输出。",
                        "先让观众知道脚会动、状态会变、代码主链在哪。",
                    ],
                ),
                SlideSpec(
                    layout="bullets",
                    title="本讲结尾要让观众带走什么",
                    bullets=[
                        "这是一个适合教学的工程型四足控制器。",
                        "先看入口和数据对象，再看控制器。",
                        "仿真入口是理解整套仓库的最快路径。",
                        "下一讲开始进入 `Controller.run()` 这个控制中枢。",
                    ],
                ),
            ],
        ),
        LectureDeck(
            slug="lecture02_controller_state_gait",
            title="第 2 讲：控制器内核、状态机与步态调度",
            subtitle="把 `Controller.run()`、状态切换和 gait 时间组织讲成观众能真正复述的骨架。",
            accent=COLOR_TEAL,
            slides=[
                SlideSpec(
                    layout="bullets",
                    title="本讲目标",
                    bullets=[
                        "看懂 `Controller.run()` 为什么是整个项目的控制中枢。",
                        "理解 `BehaviorState` 与 `activate/trot/hop` 事件关系。",
                        "看懂 `GaitController` 怎样用 tick 切分步态相位。",
                        "讲清楚对角小跑的接触表为什么能形成 trot。",
                    ],
                ),
                SlideSpec(
                    layout="bullets",
                    title="`Controller.run()` 在做什么",
                    bullets=[
                        "先根据事件更新行为状态。",
                        "根据 `REST / TROT / HOP / FINISHHOP` 进入不同控制分支。",
                        "`TROT` 分支内部再调用 `step_gait()` 做腿级路由。",
                        "最后统一落到逆运动学和 `state.joint_angles`。",
                    ],
                ),
                SlideSpec(
                    layout="two_col",
                    title="行为状态机怎么讲最清楚",
                    left_title="状态集合",
                    left_bullets=[
                        "`DEACTIVATED`",
                        "`REST`",
                        "`TROT`",
                        "`HOP`",
                        "`FINISHHOP`",
                    ],
                    right_title="事件集合",
                    right_bullets=[
                        "`activate_event`：激活/失活切换。",
                        "`trot_event`：站立与小跑切换。",
                        "`hop_event`：触发简化跳跃过程。",
                        "事件来自实机手柄层或仿真任务层。",
                    ],
                ),
                SlideSpec(
                    layout="bullets",
                    title="三条最重要的切换规则",
                    bullets=[
                        "`activate_event`：`DEACTIVATED <-> REST`。",
                        "`trot_event`：`REST <-> TROT`，也允许跳跃相关态切回 `TROT`。",
                        "`hop_event`：`REST -> HOP -> FINISHHOP -> REST`。",
                        "这一层决定宏观模式，不直接决定单腿轨迹。",
                    ],
                ),
                SlideSpec(
                    layout="code",
                    title="`TROT` 分支主链",
                    code=(
                        "state.foot_locations, contact_modes = self.step_gait(state, command)\n"
                        "rotated_foot_locations = euler2mat(command.roll, command.pitch, 0.0) @ state.foot_locations\n"
                        "(roll, pitch, yaw) = quat2euler(state.quat_orientation)\n"
                        "rmat = euler2mat(roll_compensation, pitch_compensation, 0)\n"
                        "rotated_foot_locations = rmat.T @ rotated_foot_locations\n"
                        "state.joint_angles = self.inverse_kinematics(rotated_foot_locations, self.config)\n"
                    ),
                    bullets=[
                        "先做 gait 路由，再做姿态补偿，最后做 IK。",
                        "这里是后面第 3 讲理解足端与 IK 的入口。",
                    ],
                ),
                SlideSpec(
                    layout="bullets",
                    title="`step_gait()` 的本质",
                    bullets=[
                        "每条腿单独判断当前 `contact_mode`。",
                        "`contact_mode == 1` 走 `StanceController`。",
                        "`contact_mode == 0` 走 `SwingController`。",
                        "这说明控制器是“腿级混合控制”，不是一条统一轨迹全局广播。",
                    ],
                ),
                SlideSpec(
                    layout="two_col",
                    title="步态时间基准",
                    left_title="参数",
                    left_bullets=[
                        "`dt = 0.01`",
                        "`overlap_time = 0.10`",
                        "`swing_time = 0.15`",
                        "`num_phases = 4`",
                    ],
                    right_title="推导结果",
                    right_bullets=[
                        "`overlap_ticks = 10`",
                        "`swing_ticks = 15`",
                        "`phase_ticks = [10, 15, 10, 15]`",
                        "`phase_length = 50`",
                    ],
                ),
                SlideSpec(
                    layout="code",
                    title="对角小跑的接触表",
                    code=(
                        "FR: 1 1 1 0\n"
                        "FL: 1 0 1 1\n"
                        "BR: 1 0 1 1\n"
                        "BL: 1 1 1 0\n"
                    ),
                    bullets=[
                        "`FL + BR` 一组，`FR + BL` 一组。",
                        "两组对角腿交替摆动，这就是 trot 的核心时间结构。",
                    ],
                ),
                SlideSpec(
                    layout="code",
                    title="建议现场演示命令",
                    code=(
                        "python sim/run_floating_base.py --duration 8 --mode trot --settle 1.0 --no-plots\n"
                        "python sim/run_floating_base.py --duration 8 \\\n"
                        "  --task-sequence \"rest:1.0,trot:4.0,rest\" --no-plots\n"
                    ),
                    bullets=[
                        "让观众看到行为状态和任务状态是怎样进入 `TROT` 的。",
                        "重点不是速度，而是状态切换和 gait 节律。",
                    ],
                ),
                SlideSpec(
                    layout="bullets",
                    title="本讲收束",
                    bullets=[
                        "状态机决定宏观模式。",
                        "gait 决定每条腿此刻属于 stance 还是 swing。",
                        "`ticks` 是统一时间基准。",
                        "下一讲进入 `StanceController / SwingController / IK` 细节。",
                    ],
                ),
            ],
        ),
        LectureDeck(
            slug="lecture03_foot_planning_ik",
            title="第 3 讲：足端规划、Raibert 落脚点与逆运动学",
            subtitle="把“命令 -> 足端目标 -> 关节角”这条核心链讲透。",
            accent=COLOR_GOLD,
            slides=[
                SlideSpec(
                    layout="bullets",
                    title="本讲目标",
                    bullets=[
                        "看懂 `StanceController` 和 `SwingController` 的几何意义。",
                        "理解机身速度命令如何翻译成机身坐标系下的足端运动。",
                        "讲清楚 Raibert 落脚点为什么是经典工程启发式。",
                        "讲清楚 `pupper/Kinematics.py` 怎样把足端位置变成关节角。",
                    ],
                ),
                SlideSpec(
                    layout="image",
                    title="控制器内部结构",
                    bullets=[
                        "gait scheduler 负责时间组织。",
                        "stance / swing 负责腿级足端规划。",
                        "inverse kinematics 负责把足端目标映射到关节角。",
                        "这是整套项目里最值得反复展示的一张图。",
                    ],
                    image="imgs/diagram2.jpg",
                    image_caption="控制器结构图：第 3 讲核心可视化素材。",
                ),
                SlideSpec(
                    layout="bullets",
                    title="先建立一个关键坐标直觉",
                    bullets=[
                        "足端规划大多在机身坐标系下完成。",
                        "控制器维护的是“脚相对于机身”的位置，不是先求世界坐标轨迹。",
                        "机身想向前走时，支撑足要在机身坐标里向后划地。",
                        "这是理解 stance / swing 的根本前提。",
                    ],
                ),
                SlideSpec(
                    layout="bullets",
                    title="`StanceController` 在做什么",
                    bullets=[
                        "机身想前进：足端相对机身向后移动。",
                        "机身想转向：足端绕 z 轴做反向旋转。",
                        "机身高度偏差：足端 z 方向缓慢收敛回参考高度。",
                        "所以它本质上是一个支撑相相对运动学更新器。",
                    ],
                ),
                SlideSpec(
                    layout="code",
                    title="`StanceController` 的三个关键量",
                    code=(
                        "v_xy = [-vx, -vy, (state.height - z) / z_time_constant]\n"
                        "delta_p = v_xy * dt\n"
                        "delta_R = euler2mat(0, 0, -yaw_rate * dt)\n"
                    ),
                    bullets=[
                        "`v_xy` 解释了为什么支撑足会向后划。",
                        "`delta_R` 解释了转向命令如何作用到足端。",
                        "`z_time_constant` 解释了高度误差如何被慢慢拉回。",
                    ],
                ),
                SlideSpec(
                    layout="bullets",
                    title="`SwingController` 在做什么",
                    bullets=[
                        "抬脚。",
                        "朝新的落脚点移动。",
                        "按时落下。",
                        "当前实现用三角波做 z 抬脚、用匀速逼近做 x/y 水平推进。",
                    ],
                ),
                SlideSpec(
                    layout="two_col",
                    title="Raibert 落脚点为什么重要",
                    left_title="输入影响",
                    left_bullets=[
                        "前进速度越大，落脚点越往前。",
                        "偏航速度越大，落脚点需要更多旋转补偿。",
                        "`alpha` 控制前后偏移强度。",
                        "`beta` 控制转向补偿强度。",
                    ],
                    right_title="教学价值",
                    right_bullets=[
                        "不是黑盒。",
                        "不是数值优化。",
                        "调参数时现象直观。",
                        "非常适合作为工程启发式入门。",
                    ],
                ),
                SlideSpec(
                    layout="code",
                    title="四腿逆运动学主链",
                    code=(
                        "for each leg:\n"
                        "    body_offset = config.LEG_ORIGINS[:, leg_index]\n"
                        "    local_foot = r_body_foot[:, leg_index] - body_offset\n"
                        "    joint_angles[:, leg_index] = leg_explicit_inverse_kinematics(local_foot, leg_index, config)\n"
                    ),
                    bullets=[
                        "先减去腿原点偏移，再做单腿解析逆解。",
                        "四条腿的输出统一组成 `3x4` 关节角矩阵。",
                    ],
                ),
                SlideSpec(
                    layout="bullets",
                    title="IK 里最该讲的几何量",
                    bullets=[
                        "`R_body_foot_yz`：足端在 `y-z` 平面的投影距离。",
                        "`R_hip_foot_yz`：扣除外展偏移后的有效平面距离。",
                        "`R_hip_foot`：髋到足端的空间距离。",
                        "`phi / theta / beta`：由三角几何与余弦定理得到的关节角关系。",
                    ],
                ),
                SlideSpec(
                    layout="code",
                    title="建议现场演示命令",
                    code=(
                        "python sim/run_fixed_base.py --mode rest --pitch 0.15 --roll 0.10\n"
                        "python sim/run_fixed_base.py --mode trot --x-vel 0.15 --yaw-rate 0.6\n"
                        "python sim/run_floating_base.py --duration 8 --mode trot --z-clearance 0.05 --no-plots\n"
                    ),
                    bullets=[
                        "先固定机身看几何，再浮动机身看高度变化。",
                        "第 3 讲最适合让观众第一次把现象、公式和代码对起来。",
                    ],
                ),
                SlideSpec(
                    layout="bullets",
                    title="本讲收束",
                    bullets=[
                        "stance 解决“支撑脚怎么划地”。",
                        "swing 解决“摆动脚落在哪里”。",
                        "IK 解决“足端位置怎样变成关节角”。",
                        "下一讲进入 MuJoCo 闭环、状态回灌与调参。",
                    ],
                ),
            ],
        ),
        LectureDeck(
            slug="lecture04_mujoco_closed_loop",
            title="第 4 讲：MuJoCo 闭环仿真、状态回灌与调参",
            subtitle="把 MuJoCo 讲成共享控制核心的闭环实验场，而不是孤立演示脚本。",
            accent=COLOR_BLUE,
            slides=[
                SlideSpec(
                    layout="bullets",
                    title="本讲目标",
                    bullets=[
                        "理解固定机身版和浮动机身版各自用来学什么。",
                        "看懂 `sim/sim_robot.py` 里的观测、执行器和时钟桥接。",
                        "讲清楚状态回灌、PD 力矩驱动和任务调度的意义。",
                        "给出一套真实可用的调参顺序。",
                    ],
                ),
                SlideSpec(
                    layout="two_col",
                    title="Fixed-base 与 Floating-base 的分工",
                    left_title="`run_fixed_base.py`",
                    left_bullets=[
                        "直接把关节目标写入 `qpos`。",
                        "不引入力矩和机身动力学。",
                        "最适合检查 IK、足端轨迹和姿态命令。",
                    ],
                    right_title="`run_floating_base.py`",
                    right_bullets=[
                        "机身使用 `freejoint`。",
                        "关节由 PD 力矩驱动。",
                        "姿态、速度、触地和关节状态都会回写到 `State`。",
                    ],
                ),
                SlideSpec(
                    layout="image",
                    title="浮动机身控制流",
                    bullets=[
                        "任务层负责生成 `Command`。",
                        "控制器核心保持不变。",
                        "执行器层把关节目标角转换成 PD 力矩。",
                        "MuJoCo 产生命令响应与观测，再回写到 `State`。",
                    ],
                    image="docs/mujoco_quadruped_mastery_tutorial/assets/floating_base_control_flow.png",
                    image_caption="浮动机身版：最适合讲清闭环主链的一张图。",
                ),
                SlideSpec(
                    layout="bullets",
                    title="`sim_robot.py` 的五个桥接模块",
                    bullets=[
                        "`SimObservationInterface`：传感器桥接。",
                        "`SimIMU`：仿真版 IMU 适配器。",
                        "`SimHardwareInterface`：关节目标角到 PD 力矩。",
                        "`SimControlClock`：控制周期和仿真子步对齐。",
                        "`TaskCommandSource`：无手柄条件下的任务级命令源。",
                    ],
                ),
                SlideSpec(
                    layout="image",
                    title="为什么说这是闭环",
                    bullets=[
                        "姿态会回写。",
                        "机身速度会回写。",
                        "触地状态会回写。",
                        "测量足端位置还会融合回 `state.foot_locations`。",
                    ],
                    image="docs/mujoco_quadruped_mastery_tutorial/assets/state_feedback_closed_loop.png",
                    image_caption="状态反馈图：观众第一次真正理解闭环的关键素材。",
                ),
                SlideSpec(
                    layout="code",
                    title="`sync_state()` 的关键职责",
                    code=(
                        "state.measured_foot_locations = measured_feet\n"
                        "state.measured_joint_angles = measured_joint_angles\n"
                        "state.body_position = torso_position\n"
                        "state.body_velocity = imu_vel\n"
                        "state.angular_velocity = imu_gyro\n"
                        "state.contact_estimate = foot_forces > contact_threshold\n"
                    ),
                    bullets=[
                        "这一层让仿真里的机器人“长得像”真实机器人的观测接口。",
                        "也让后续姿态反馈、触地估计和日志观察成为可能。",
                    ],
                ),
                SlideSpec(
                    layout="bullets",
                    title="PD 力矩桥接怎么讲",
                    bullets=[
                        "控制器输出仍然是目标关节角。",
                        "MuJoCo 执行器需要的是 motor 控制量。",
                        "`SimHardwareInterface` 用 `kp / kd / torque_limit` 做关节级跟踪。",
                        "因此这里最适合讲“刚度、阻尼、饱和”三件事。",
                    ],
                ),
                SlideSpec(
                    layout="code",
                    title="任务序列语法",
                    code=(
                        "mode[:duration][@key=value;key=value...],...\n\n"
                        "python sim/run_floating_base.py --duration 8 \\\n"
                        "  --task-sequence \"rest:1.0,trot:4.0@vx=0.08;z_clearance=0.04,rest\" \\\n"
                        "  --no-plots\n"
                    ),
                    bullets=[
                        "没有手柄时，也能做稳定、可重复的实验录制。",
                        "非常适合视频课里做分段演示和参数对比。",
                    ],
                ),
                SlideSpec(
                    layout="two_col",
                    title="推荐调参顺序",
                    left_title="先做什么",
                    left_bullets=[
                        "先用 `rest` 确认站姿和初始高度。",
                        "再用 `fixed-base` 检查 IK 和足端轨迹。",
                        "再切到 `floating-base` 小速度前进。",
                    ],
                    right_title="再做什么",
                    right_bullets=[
                        "调 `kp / kd / torque-limit`。",
                        "调 `overlap_time / swing_time / z_clearance`。",
                        "最后再引入 `attitude_kp / attitude_kd / velocity_kp`。",
                    ],
                ),
                SlideSpec(
                    layout="bullets",
                    title="本讲收束",
                    bullets=[
                        "fixed-base 用来看几何，floating-base 用来看闭环。",
                        "`sim_robot.py` 是共享控制核心与 MuJoCo 之间的关键桥。",
                        "第 4 讲讲清楚后，观众就有能力自己做实验与调参。",
                        "下一讲进入实机接口、Sim2Real 与扩展路线。",
                    ],
                ),
            ],
        ),
        LectureDeck(
            slug="lecture05_hardware_sim2real",
            title="第 5 讲：实机接口、Sim2Real 与二次开发路线",
            subtitle="把实机链路、仿真到实机的边界和后续扩展路线讲清楚。",
            accent=COLOR_TEAL,
            slides=[
                SlideSpec(
                    layout="bullets",
                    title="本讲目标",
                    bullets=[
                        "看懂实机运行链里哪些层和仿真不同。",
                        "理解 `JoystickInterface`、`HardwareInterface`、IMU 和部署脚本的角色。",
                        "总结从仿真迁移到实机时最容易踩的坑。",
                        "给出一个适合学生继续做下去的扩展路线。",
                    ],
                ),
                SlideSpec(
                    layout="code",
                    title="实机控制链",
                    code=(
                        "PS4 手柄\n"
                        "  -> UDP\n"
                        "  -> JoystickInterface\n"
                        "  -> Command\n"
                        "  -> Controller\n"
                        "  -> joint angles\n"
                        "  -> HardwareInterface\n"
                        "  -> pigpio PWM\n"
                        "  -> servos\n"
                    ),
                    bullets=[
                        "这条链和仿真版最大的区别，在于命令源、执行器和观测后端。",
                        "共享的仍然是 `Controller`、gait、stance、swing 和 IK。",
                    ],
                ),
                SlideSpec(
                    layout="bullets",
                    title="`JoystickInterface` 的价值",
                    bullets=[
                        "负责把 UDP 手柄消息变成统一的 `Command`。",
                        "连续量包括 `horizontal_velocity / yaw_rate / pitch / height / roll`。",
                        "离散按钮采用边沿触发，不是按下即持续重复触发。",
                        "pitch 还会经过死区和一阶限速滤波，这就是命令整形。",
                    ],
                ),
                SlideSpec(
                    layout="bullets",
                    title="`HardwareInterface` 的价值",
                    bullets=[
                        "连接本机 `pigpio` 守护进程。",
                        "维护 PWM 引脚、频率和输出范围。",
                        "把关节角转换成 PWM 脉宽和 duty cycle。",
                        "把数学模型里的角度真正送到实机舵机。",
                    ],
                ),
                SlideSpec(
                    layout="two_col",
                    title="为什么标定不能省",
                    left_title="需要标定的原因",
                    left_bullets=[
                        "模型零位不等于机械零位。",
                        "不同腿和不同关节的方向不一致。",
                        "`servo_multipliers` 决定角度映射方向。",
                    ],
                    right_title="课程里该怎么讲",
                    right_bullets=[
                        "没有标定，再好的 gait 也站不稳。",
                        "标定是 Sim2Real 里最现实的一关。",
                        "它体现了工程系统不是只有算法，还有装配与偏差。",
                    ],
                ),
                SlideSpec(
                    layout="bullets",
                    title="IMU 与轻量闭环",
                    bullets=[
                        "IMU 在 `run_robot.py` 中是可选的，但非常重要。",
                        "它说明原始实机入口就预留了姿态反馈。",
                        "这不是完整状态估计器，但已经不是纯几何开环。",
                        "很适合作为从开环 gait 走向轻量闭环的过渡案例。",
                    ],
                ),
                SlideSpec(
                    layout="image",
                    title="从仿真到实机的迁移边界",
                    bullets=[
                        "共享：`Controller`、gait、stance、swing、IK。",
                        "替换：命令源、执行器接口、观测接口。",
                        "这张图很适合用来收束整套系列课。",
                    ],
                    image="docs/mujoco_quadruped_mastery_tutorial/assets/sim_to_real_migration.png",
                    image_caption="Sim2Real：共享控制核心，替换平台相关后端。",
                ),
                SlideSpec(
                    layout="bullets",
                    title="实机最容易踩的坑",
                    bullets=[
                        "舵机中位和机械装配偏差。",
                        "仿真里的 `kp / kd` 不能直接等价搬到真实舵机。",
                        "地面摩擦、接触和机身重量分布都会改变。",
                        "手柄输入是异步网络消息，和本地任务源不同。",
                    ],
                ),
                SlideSpec(
                    layout="two_col",
                    title="推荐的二次开发路线",
                    left_title="先做工程清理",
                    left_bullets=[
                        "给 `run_robot.py` 增加 `__main__` 保护。",
                        "改进主循环定时方式。",
                        "统一日志与观测输出。",
                    ],
                    right_title="再做能力增强",
                    right_bullets=[
                        "替换摆腿轨迹。",
                        "增加新步态表和姿态反馈。",
                        "继续走向阻抗、WBC、MPC 或学习方法。",
                    ],
                ),
                SlideSpec(
                    layout="bullets",
                    title="整套课的最终目标",
                    bullets=[
                        "不是背术语，而是讲清主链。",
                        "不是只会跑一个脚本，而是知道如何扩展。",
                        "不是停在仿真，而是知道如何走向实机。",
                        "这一讲结束后，观众应当能自己规划下一步项目路线。",
                    ],
                ),
            ],
        ),
    ]


def main() -> None:
    output_paths = [create_deck(deck) for deck in build_lectures()]
    for output_path in output_paths:
        print(output_path)


if __name__ == "__main__":
    main()
